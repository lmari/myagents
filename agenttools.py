from openai_functools import openai_function
import json

def _get_tool_metadata(tool) -> dict:
    """ Genera i metadati per la funzione specificata. """
    return {'type': 'function', 'function': tool.openai_metadata }

def _exec_tool(response) -> dict:
    """
    Se la risposta dal LM contiene l'informazione per chiamare una o più funzioni,
    esegue una per una queste funzioni e restituisce un dizionario con 'correct'=True, il risultato e altre informazioni
    altrimenti restituisce un dizionario con 'correct'=False e un messaggio di errore come risultato.
    """
    result = {"correct": True}
    if not response.choices[0].message.tool_calls:
        result["correct"] = False
        result["result"] = "No function call found in the response." # type: ignore
        return result
    function_call_results = []
    for_completion_messages = []
    for tool_call in response.choices[0].message.tool_calls:
        function_call_result = {}
        for_completion_message = {"role": "tool"}
        try:
            function_name = tool_call.function.name
            function_arguments = json.loads(tool_call.function.arguments)
            function_call_result["function_name"] = function_name
            function_call_result["function_arguments"] = function_arguments
            for_completion_message["tool_call_id"] = tool_call.id
        except Exception as e:
            result["correct"] = False
            result["result"] = f"Error parsing function call: {e}" # type: ignore
            return result
        if function_name not in globals():
            result["correct"] = False
            result["result"] = f"Function {function_name} not found." # type: ignore
            return result
        try:
            function_call_result["result"] = str(globals()[function_name](**function_arguments))
            function_call_results.append(function_call_result)
            for_completion_message["content"] = function_call_result["result"]
            for_completion_messages.append(for_completion_message)
        except Exception as e:
            result["correct"] = False
            result["result"] = f"Error executing function {function_name}: {e}" # type: ignore
            return result
    result["result"] = function_call_results # type: ignore
    result["for_completion_messages"] = for_completion_messages # type: ignore
    return result


# *********************
# Exemplary classes ***
# *********************

import os

@openai_function
def list_files(path: str) -> list | None:
    """
    Elenca i file contenuti nella directory specificata, o restituisce None se la directory non esiste.
    Args:
        path: Il nome della directory di cui elencare i file
    Returns:
        L'elenco dei file nella directory, o None se la directory non esiste
    """
    path = path.strip()
    if not os.path.isdir(path):
        return None
    return os.listdir(path)

@openai_function
def read_file(path: str) -> str | None:
    """
    Legge il contenuto di un file e lo restituisce come stringa, o restituisce None se il file non esiste.
    Args:
        path: Il nome del file da leggere
    Returns:
        Il contenuto del file come stringa, o None se il file non esiste
    """
    path = path.strip()
    if not os.path.isfile(path):
        return None
    with open(path, 'r') as f:
        return f.read()
    
# ***************************************************************************

from duckduckgo_search import DDGS

@openai_function
def web_search(query: str, max_results: int=3) -> list:
    """
    Cerca nel web pagine che abbiano come argomento la condizione di ricerca specificata.
    Args:
        query: La condizione di ricerca
        max_results: Il numero massimo di risultati da restituire (3 se non specificato altrimenti)
    Returns:
        list: L'elenco di risultati della ricerca, o None se non sono stati trovati risultati
    """
    return DDGS().text(query, max_results=max_results)

# ***************************************************************************

@openai_function
def evaluate_expression(code: str) -> str | None:
    """
    Valuta l'espressione matematica specificata e restituisce il risultato, o restituisce il messaggio di errore generato.
    Args:
        code: Il codice Python da valutare
    Returns:
        Il risultato della valutazione dell'espressione o un messaggio di errore
    """
    code = code.strip()
    if not code:
        return None
    try:
        return str(eval(code))
    except Exception as e:
        return str(e)

# ***************************************************************************

import random

@openai_function
def generate_random_number(min_value: float, max_value: float) -> float:
    """
    Restituisce un numero casuale tra il valore minimo e il valore massimo specificati.
    Args:
        min_value: Il valore minimo
        max_value: Il valore massimo
    Returns:
        Un numero casuale tra il valore minimo e il valore massimo
    """
    return random.uniform(min_value, max_value)

# ***************************************************************************

from datetime import datetime
import pytz

@openai_function
def get_current_day_and_time(zone: str="Europe/Rome") -> datetime:
    """
    Restituisce la data e l'orario corrente nella zona specificata, o in Italia se non specificato.
    Args:
        zone: La zona di interesse (Italiana se non specificata)
    Returns:
        La data e l'orario corrente
    """
    return datetime.now(pytz.timezone(zone))

# ***************************************************************************

from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import markdown, re
from bs4 import BeautifulSoup

@openai_function
def markdown_to_pptx(md_text: str, output_path: str="slides.pptx") -> str:
    """
    Genera una presentazione PowerPoint da un testo in formato Markdown.
    Args:
        md_text: Il testo in formato Markdown da cui generare la presentazione
        output_path: Il percorso del file di output (default "slides.pptx")
    Returns:
        Il percorso del file con la presentazione generata
    """
    md_text = md_text[md_text.find("#"):] # remove from md_text everything before the first # (title)
    #md_text = md_text[:md_text.rfind("\n")] # remove from md_text everything after the last \n included

    prs = Presentation()
    layout = prs.slide_layouts[1] # titolo + contenuto
    chunks = re.split(r'(^# .+$)', md_text, flags=re.MULTILINE)
    chunks = [c for c in chunks if c.strip()]   # elimina stringhe vuote
    for i in range(0, len(chunks), 2):
        title_line = chunks[i]
        body_md = chunks[i+1] if i+1 < len(chunks) else ""
        slide = prs.slides.add_slide(layout)
        titleph = slide.shapes.title
        bodyph = slide.shapes.placeholders[1].text_frame # type: ignore
        titleph.text = title_line.lstrip("# ").strip() # type: ignore
        bodyph.clear()
        html = markdown.markdown(body_md)
        soup = BeautifulSoup(html, "html.parser")

        def add_run(paragraph, node):
            run = paragraph.add_run()
            run.text = node.get_text() if hasattr(node, "get_text") else str(node)
            if getattr(node, "name", None) in ("strong", "b"):
                run.font.bold = True
            if getattr(node, "name", None) in ("em", "i"):
                run.font.italic = True
        
        for elem in soup.children:
            if elem.name in ("p", "ul", "ol"): # type: ignore
                if elem.name == "p": # type: ignore
                    p = bodyph.add_paragraph()
                    p.level = 0
                    for child in elem.children: # type: ignore
                        add_run(p, child if hasattr(child, "name") else elem)
                else:
                    for li in elem.find_all("li", recursive=False): # type: ignore
                        p = bodyph.add_paragraph()
                        p.level = 1
                        for child in li.children:
                            add_run(p, child if hasattr(child, "name") else li)
            elif elem.name == "h2": # type: ignore
                p = bodyph.add_paragraph()
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                p.level = 0
                add_run(p, elem)

    prs.save(output_path)
    return output_path
